#!/usr/bin/env python3
"""生成 HDB 组屋转售价格分析答辩 PPT — 更新版（6 页面系统）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ============================================================
# 配色方案 — 新加坡 HDB 蓝白主题
# ============================================================
DARK_BLUE   = RGBColor(0x0B, 0x2D, 0x5B)
PRIMARY_BLUE = RGBColor(0x1A, 0x56, 0xDB)
LIGHT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
ACCENT_RED  = RGBColor(0xEF, 0x44, 0x44)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1E, 0x29, 0x3B)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY    = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

CENTER = PP_ALIGN.CENTER
LEFT   = PP_ALIGN.LEFT
RIGHT  = PP_ALIGN.RIGHT

# ============================================================
# 工具函数
# ============================================================
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=BLACK, bold=False, alignment=LEFT, font_name="Microsoft YaHei",
                line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * font_size - font_size)
    return tf

def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, font_size=16, color=BLACK, bold=False, alignment=LEFT,
             font_name="Microsoft YaHei", space_after=6, space_before=0,
             level=0, bullet=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    return p

def add_bottom_bar(slide):
    add_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), PRIMARY_BLUE)

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.3), Inches(7.12), Inches(0.8), Inches(0.35),
                str(num), font_size=10, color=WHITE, alignment=CENTER)

def slide_title(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.7),
                title, font_size=36, color=DARK_BLUE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.45),
                    subtitle, font_size=16, color=GRAY)
    add_rect(slide, Inches(0.8), Inches(1.55), Inches(1.5), Inches(0.05), ACCENT_BLUE)

def make_table(slide, left, top, col_widths, headers, rows, header_color=None):
    if header_color is None:
        header_color = DARK_BLUE
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_width = sum(col_widths)
    row_h = Inches(0.42)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_width,
                                          row_h * n_rows)
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Microsoft YaHei"
            p.alignment = CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BLACK
                p.font.name = "Microsoft YaHei"
                p.alignment = CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape

def make_card(slide, left, top, width, height, title, value, color=PRIMARY_BLUE):
    card = add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, width, Inches(0.06), color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.45),
                value, font_size=24, color=color, bold=True, alignment=LEFT)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.58), width - Inches(0.4), Inches(0.3),
                title, font_size=11, color=GRAY, alignment=LEFT)


# ============================================================
# 第 1 页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), ACCENT_BLUE)
add_rect(slide, Inches(0.8), Inches(2.2), Inches(0.08), Inches(3.2), ACCENT_BLUE)

add_textbox(slide, Inches(1.3), Inches(2.2), Inches(11), Inches(1.0),
            "新加坡 HDB 组屋转售价格", font_size=48, color=WHITE, bold=True)
add_textbox(slide, Inches(1.3), Inches(3.1), Inches(11), Inches(0.8),
            "分析与预测系统", font_size=48, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(1.3), Inches(4.1), Inches(10.5), Inches(0.6),
            "基于 Streamlit 的交互式房产数据分析平台", font_size=22, color=RGBColor(0x94, 0xA3, 0xB8))
add_rect(slide, Inches(1.3), Inches(4.9), Inches(3.5), Inches(0.03), ACCENT_BLUE)
add_textbox(slide, Inches(1.3), Inches(5.2), Inches(5), Inches(0.45),
            "课程项目答辩", font_size=18, color=WHITE)
add_textbox(slide, Inches(1.3), Inches(5.65), Inches(5), Inches(0.45),
            "小组成员：张涛、申经纬", font_size=16, color=RGBColor(0x94, 0xA3, 0xB8))
add_textbox(slide, Inches(1.3), Inches(6.1), Inches(5), Inches(0.45),
            "2026 年 6 月", font_size=16, color=GRAY)


# ============================================================
# 第 2 页：项目背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 2)
slide_title(slide, "项目背景", "Why this project?")

tf = add_rich_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.4))
add_para(tf, "📊  为什么要做这个项目？", font_size=20, color=DARK_BLUE, bold=True, space_after=12)
add_para(tf, "新加坡超 80% 人口居住在 HDB 组屋", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "转售市场活跃，房价受多重因素影响", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "购房者面对海量数据，缺乏直观分析工具", font_size=15, color=BLACK, space_after=16, level=1)

add_para(tf, "🎯  我们要解决什么？", font_size=20, color=DARK_BLUE, bold=True, space_after=12)
add_para(tf, "如何快速了解某个区域的房价全貌？", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "哪些因素对房价影响最大？", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "能否用模型预测价格？预算有限买哪里最划算？", font_size=15, color=BLACK, space_after=6, level=1)

card = add_rect(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.0), LIGHT_GRAY)
add_rect(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.08), PRIMARY_BLUE)
tf2 = add_rich_textbox(slide, Inches(7.4), Inches(2.3), Inches(4.7), Inches(3.3))
add_para(tf2, "📌  研究对象", font_size=20, color=DARK_BLUE, bold=True, space_after=16)
add_para(tf2, "🏘️  3 个相邻镇区", font_size=16, color=BLACK, bold=True, space_after=4)
add_para(tf2, "Punggol（榜鹅）· Sengkang（盛港）· Hougang（后港）", font_size=14, color=GRAY, space_after=12)
add_para(tf2, "📡  数据来源", font_size=16, color=BLACK, bold=True, space_after=4)
add_para(tf2, "data.gov.sg 官方 HDB 转售交易 API", font_size=14, color=GRAY, space_after=12)
add_para(tf2, "📅  时间范围", font_size=16, color=BLACK, bold=True, space_after=4)
add_para(tf2, "2020 年 — 2026 年", font_size=14, color=GRAY, space_after=12)
add_para(tf2, "📋  数据规模", font_size=16, color=BLACK, bold=True, space_after=4)
add_para(tf2, "34,914 条清洗后交易记录", font_size=18, color=ACCENT_RED, bold=True, space_after=4)


# ============================================================
# 第 3 页：系统架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 3)
slide_title(slide, "系统架构", "System Architecture")

# Streamlit Web App 外壳
add_rect(slide, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.2), LIGHT_GRAY)
add_textbox(slide, Inches(0.6), Inches(2.15), Inches(12.1), Inches(0.45),
            "Streamlit Web App", font_size=18, color=WHITE, bold=True, alignment=CENTER,
            font_name="Consolas")

# 6 个页面模块
modules = ["📊 数据概览", "🗺️ 地图可视化", "📈 因素分析", "🤖 价格预测", "💭 分析思考", "🏆 策略验证"]
mod_width = Inches(1.78)
start_x = Inches(0.75)
gap = Inches(0.18)
for i, mod in enumerate(modules):
    x = start_x + i * (mod_width + gap)
    add_rect(slide, x, Inches(2.65), mod_width, Inches(0.55), PRIMARY_BLUE)
    add_textbox(slide, x, Inches(2.68), mod_width, Inches(0.5),
                mod, font_size=11, color=WHITE, bold=True, alignment=CENTER)

# 数据管线层
add_rect(slide, Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.7), DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(3.55), Inches(12.1), Inches(0.6),
            "data_loader.py  数据管线：data.gov.sg API → 清洗 → Parquet 缓存",
            font_size=15, color=WHITE, bold=True, alignment=CENTER, font_name="Consolas")

# 技术栈
tf = add_rich_textbox(slide, Inches(0.6), Inches(4.6), Inches(5.8), Inches(2.0))
add_para(tf, "🖥️  前端技术", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "Streamlit + Plotly + pydeck", font_size=14, color=BLACK, space_after=14)
add_para(tf, "⚙️  后端处理", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "Pandas + NumPy + Scikit-learn", font_size=14, color=BLACK, space_after=14)
add_para(tf, "📈  统计分析", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "SciPy + Statsmodels", font_size=14, color=BLACK, space_after=14)

tf2 = add_rich_textbox(slide, Inches(7.0), Inches(4.6), Inches(5.8), Inches(2.0))
add_para(tf2, "🚀  部署", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf2, "Streamlit Community Cloud", font_size=14, color=BLACK, space_after=14)
add_para(tf2, "📦  数据格式", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf2, "Parquet（高性能列式存储）", font_size=14, color=BLACK, space_after=14)
add_para(tf2, "🔑  关键特性", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf2, "一键部署 · 内置缓存 · API 降级", font_size=14, color=BLACK, space_after=14)


# ============================================================
# 第 4 页：数据概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 4)
slide_title(slide, "数据概览 — 功能展示", "Data Overview")

tf = add_rich_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.5))
add_para(tf, "🔎  多条件筛选", font_size=20, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "镇区 / 年份 / 房型 / 面积 / 价格 / 剩余年限 / 楼层", font_size=14, color=GRAY, space_after=16)

add_para(tf, "📊  统计看板（8 指标卡片）", font_size=20, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "交易套数 · 均价 · 中位数 · 总成交额", font_size=14, color=BLACK, space_after=4)
add_para(tf, "最高价 · 最低价 · 最高单价 · 最低单价", font_size=14, color=BLACK, space_after=4)
add_para(tf, "镇区分组统计：套数、均价、中位数、单价、成交总额", font_size=14, color=GRAY, space_after=16)

add_para(tf, "📈  趋势 + 分布", font_size=20, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "月度均价走势（分镇区折线图）", font_size=14, color=BLACK, space_after=4)
add_para(tf, "季度交易量（分镇区柱状图）", font_size=14, color=BLACK, space_after=4)
add_para(tf, "价格分布直方图", font_size=14, color=BLACK, space_after=4)
add_para(tf, "镇区 × 房型均价热力图", font_size=14, color=BLACK, space_after=16)

add_para(tf, "📋  交互式数据表", font_size=20, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "可排序 · 可筛选 · 可导出 CSV", font_size=14, color=BLACK, space_after=4)

# 右栏 — 模拟统计卡片
cards_data = [
    ("交易套数", "34,914", PRIMARY_BLUE),
    ("均价", "S$ 558K", ACCENT_GREEN),
    ("中位数", "S$ 545K", ACCENT_BLUE),
    ("总成交额", "S$ 19.5B", ACCENT_ORANGE),
    ("最高价", "S$ 1.38M", ACCENT_RED),
    ("最低价", "S$ 228K", DARK_BLUE),
]
for i, (label, val, clr) in enumerate(cards_data):
    row = i // 3
    col = i % 3
    x = Inches(7.2) + col * Inches(2.0)
    y = Inches(1.9) + row * Inches(1.1)
    make_card(slide, x, y, Inches(1.85), Inches(0.9), label, val, clr)


# ============================================================
# 第 5 页：地图可视化
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 5)
slide_title(slide, "地图可视化 — 三种模式", "Map Visualization — pydeck")

make_table(slide, Inches(0.8), Inches(1.9),
           [Inches(1.6), Inches(3.6), Inches(3.2), Inches(2.8)],
           ["模式", "功能", "技术", "交互"],
           [
               ["镇区均价", "ScatterplotLayer\n大小=交易量，颜色=均价", "pydeck + 全局色阶", "悬停详情弹出"],
               ["热力分布", "房价热力图\n含坐标抖动仿真（~300m）", "pydeck HeatmapLayer", "缩放/平移/强度调节"],
               ["时空变化", "年份滑块\n统一色阶跨年对比", "pydeck + Plotly\n趋势图联动", "时间维度切换"],
           ])

tf = add_rich_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.3))
add_para(tf, "🔌  配套设施叠加（已实现）", font_size=20, color=DARK_BLUE, bold=True, space_after=12)
add_para(tf, "🚇  MRT / LRT 站点 — NE14-NE17 + Punggol/Sengkang LRT 线路", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "🏫  学校 — 5 所小学 + 5 所中学（Punggol/Sengkang/Hougang 区域）", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "🎛️  侧边栏多选叠加控制 + 热力半径/强度滑块调节", font_size=15, color=GRAY, space_after=6, level=1)


# ============================================================
# 第 6 页：影响因素分析 — 7 个 Tab
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 6)
slide_title(slide, "影响因素分析 — 7 个分析 Tab", "Factor Analysis")

make_table(slide, Inches(0.25), Inches(1.85),
           [Inches(1.5), Inches(2.2), Inches(4.0), Inches(4.3)],
           ["Tab", "方法", "发现", "统计显著性"],
           [
               ["① 面积", "散点图 + OLS + Pearson r", "面积核心驱动单价，r≈0.75", "★★★★★"],
               ["② 房型", "箱线图 + ANOVA", "p<0.001，每级溢价 15-25%", "★★★★★"],
               ["③ 楼层", "箱线图 + 趋势分析", "中高层溢价 5-10%", "★★★"],
               ["④ 剩余租约", "分箱箱线图 + 线性回归", "租约越短折价越大", "★★★★"],
               ["⑤ 房龄", "散点图 + LOWESS 趋势", "新房显著溢价", "★★★★"],
               ["⑥ 镇区对比", "Town vs Price 偏离度 + CAGR", "Hougang 均价最高", "★★★"],
               ["💡 策略验证", "4 策略 CAGR + CV + 雷达图", "各策略适用场景分析", "—"],
           ])

# 结论
card = add_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), LIGHT_GRAY)
add_rect(slide, Inches(0.5), Inches(5.5), Inches(0.08), Inches(1.2), ACCENT_GREEN)
tf = add_rich_textbox(slide, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.9))
add_para(tf, "📊  核心结论", font_size=20, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "面积是最强单一预测因子（r≈0.75），房型与面积高度相关存在共线性；镇区对比采用 Town vs 均价偏离度框架，Hougang 均价最高、Punggol 增长最快",
         font_size=14, color=BLACK, space_after=4)


# ============================================================
# 第 7 页：价格预测 — 双模型对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 7)
slide_title(slide, "价格预测 — 双模型对比", "Price Prediction: Ridge vs Random Forest")

# 说明框
tf0 = add_rich_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.6))
add_para(tf0, "预测目标：单价 price_per_sqm（新元/sqm） | 时间切分：2020–2023 训练 / 2024+ 测试 | 特征：8 个", font_size=15, color=GRAY)

make_table(slide, Inches(0.8), Inches(2.4),
           [Inches(2.5), Inches(2.8), Inches(3.0), Inches(3.0)],
           ["指标", "Ridge Regression", "Random Forest", "说明"],
           [
               ["类型", "线性（L2 正则）", "集成树（250 棵树）", "—"],
               ["MAE (新元/sqm)", "~S$520/sqm ✅", "~S$700/sqm", "越低越好"],
               ["R²", "~0.48 ✅", "~0.19", "越高越好"],
               ["MAPE", "~7.7% ✅", "~10.0%", "越低越好"],
               ["5-fold CV R²", "稳定 ✅", "泛化有限", "时间外推能力"],
           ])

# 结论卡片
card = add_rect(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(1.1), LIGHT_GRAY)
add_rect(slide, Inches(0.8), Inches(4.3), Inches(0.08), Inches(1.1), ACCENT_GREEN)
tf = add_rich_textbox(slide, Inches(1.2), Inches(4.4), Inches(10.8), Inches(0.9))
add_para(tf, "🏆  Ridge 在时间切分验证中表现更优", font_size=22, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "原因：Ridge 的 year 特征可线性外推价格时间趋势（2020-2023 → 2024+）；RF 树模型无法外推时间趋势，在纯时间切分中泛化受限",
         font_size=14, color=BLACK, space_after=6)

# 诊断 + 分房型评估
tf2 = add_rich_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.0))
add_para(tf2, "🔍  模型诊断 & 分房型评估", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf2, "6 类别分房型评估（小户型/中户型/大户型/老旧/新近/MRT沿线）· 预测 vs 实际散点图 · 残差分布诊断 · 5-fold CV",
         font_size=14, color=BLACK, space_after=4)


# ============================================================
# 第 8 页：特征重要性
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 8)
slide_title(slide, "特征重要性", "Feature Importance")

make_table(slide, Inches(0.5), Inches(1.9),
           [Inches(1.2), Inches(3.5), Inches(3.0), Inches(3.8)],
           ["排名", "特征", "重要性", "解读"],
           [
               ["🥇 1", "floor_area_sqm（面积）", "★★★★★", "绝对领先，最强预测因子"],
               ["🥈 2", "flat_type_ordinal（房型编码）", "★★★★", "房型等级与价格强相关"],
               ["🥉 3", "remaining_lease（剩余租约）", "★★★", "租约越短折价越大"],
               ["4", "year（交易年份）", "★★★", "市场整体涨跌趋势"],
               ["5", "storey_mid（楼层）", "★★", "中高层溢价明显"],
               ["6", "town（镇区）", "★★", "地段差异被面积部分解释"],
               ["7", "mrt_distance（MRT 距离）", "★", "步行范围内影响有限"],
               ["8", "schools_nearby（学校数量）", "★", "学区效应较弱"],
           ])

card = add_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3), LIGHT_GRAY)
add_rect(slide, Inches(0.5), Inches(5.3), Inches(0.08), Inches(1.3), ACCENT_ORANGE)
tf = add_rich_textbox(slide, Inches(0.9), Inches(5.45), Inches(11.5), Inches(1.0))
add_para(tf, "💡  关键发现", font_size=20, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "面积仍是最强预测因子；房型/镇区特征的重要性被面积变量部分替代（面积和房型高度相关）", font_size=15, color=BLACK, space_after=4)
add_para(tf, "Ridge 标准化系数排名与 RF 重要性排名高度一致，验证了结论的可靠性", font_size=15, color=GRAY, space_after=4)


# ============================================================
# 第 9 页：交互式预测
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 9)
slide_title(slide, "交互式预测", "Interactive Prediction")

# 输入区
add_rect(slide, Inches(0.8), Inches(2.1), Inches(3.5), Inches(3.8), LIGHT_GRAY)
add_rect(slide, Inches(0.8), Inches(2.1), Inches(3.5), Inches(0.08), PRIMARY_BLUE)
tf = add_rich_textbox(slide, Inches(1.1), Inches(2.35), Inches(3.0), Inches(3.3))
add_para(tf, "📝  用户输入", font_size=18, color=DARK_BLUE, bold=True, space_after=12)
add_para(tf, "镇区（3 选 1）", font_size=14, color=BLACK, space_after=4)
add_para(tf, "房型（6 种）", font_size=14, color=BLACK, space_after=4)
add_para(tf, "面积（㎡）", font_size=14, color=BLACK, space_after=4)
add_para(tf, "楼层范围", font_size=14, color=BLACK, space_after=4)
add_para(tf, "剩余年限", font_size=14, color=BLACK, space_after=4)
add_para(tf, "交易年份", font_size=14, color=BLACK, space_after=4)

add_textbox(slide, Inches(4.5), Inches(3.2), Inches(1.2), Inches(0.5),
            "➡️ 模型 ➡️", font_size=16, color=PRIMARY_BLUE, bold=True, alignment=CENTER)

# 输出区
add_rect(slide, Inches(5.9), Inches(2.1), Inches(6.6), Inches(3.8), LIGHT_GRAY)
add_rect(slide, Inches(5.9), Inches(2.1), Inches(6.6), Inches(0.08), ACCENT_GREEN)
tf2 = add_rich_textbox(slide, Inches(6.2), Inches(2.35), Inches(6.0), Inches(3.3))
add_para(tf2, "📤  预测输出", font_size=18, color=DARK_BLUE, bold=True, space_after=12)
add_para(tf2, "Ridge 预测: S$XXX,XXX ⭐ 推荐", font_size=16, color=ACCENT_GREEN, bold=True, space_after=4)
add_para(tf2, "RF 预测: S$XXX,XXX", font_size=14, color=GRAY, space_after=8)
add_para(tf2, "📐  90% 预测区间（250 棵个体树分位数）", font_size=14, color=DARK_BLUE, bold=True, space_after=4)
add_para(tf2, "[S$XXX,XXX  —  S$XXX,XXX]", font_size=14, color=BLACK, space_after=8)
add_para(tf2, "📋  历史相似成交（渐进式放宽匹配）", font_size=14, color=DARK_BLUE, bold=True, space_after=4)
add_para(tf2, "4 级渐进放宽：面积±10→±30sqm + 租约放宽 + 兜底策略", font_size=14, color=GRAY, space_after=4)


# ============================================================
# 第 10 页：分析思考题（NEW）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 10)
slide_title(slide, "分析思考题", "Analysis & Thinking Questions — 3 Questions × 200+ Words")

# Q1
add_rect(slide, Inches(0.5), Inches(1.9), Inches(3.9), Inches(2.5), LIGHT_GRAY)
add_rect(slide, Inches(0.5), Inches(1.9), Inches(3.9), Inches(0.06), PRIMARY_BLUE)
tf1 = add_rich_textbox(slide, Inches(0.75), Inches(2.1), Inches(3.4), Inches(2.2))
add_para(tf1, "Q1：HDB 转售保值性", font_size=17, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf1, "面积—价格基本盘 → 剩余租约—时间折扣 → 镇区—地段溢价", font_size=13, color=BLACK, space_after=6)
add_para(tf1, "CAGR 年化复合增长率量化 + 面积/房龄/镇区分组对比分析", font_size=13, color=GRAY, space_after=6)
add_para(tf1, "✅ 含 Haversine MRT 距离计算 + 200+ 字分析", font_size=12, color=ACCENT_GREEN, space_after=4)

# Q2
add_rect(slide, Inches(4.7), Inches(1.9), Inches(3.9), Inches(2.5), LIGHT_GRAY)
add_rect(slide, Inches(4.7), Inches(1.9), Inches(3.9), Inches(0.06), ACCENT_ORANGE)
tf2 = add_rich_textbox(slide, Inches(4.95), Inches(2.1), Inches(3.4), Inches(2.2))
add_para(tf2, "Q2：购房策略方向论证", font_size=17, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf2, "4 策略逻辑推演：镇区偏离回归 / MRT 便利 / 长剩余租约 / 低总价入门", font_size=13, color=BLACK, space_after=6)
add_para(tf2, "策略对应数据逻辑 + Train/Test 验证支撑", font_size=13, color=GRAY, space_after=6)
add_para(tf2, "✅ 只用 3 镇区数据 + 200+ 字分析", font_size=12, color=ACCENT_GREEN, space_after=4)

# Q3
add_rect(slide, Inches(8.9), Inches(1.9), Inches(3.9), Inches(2.5), LIGHT_GRAY)
add_rect(slide, Inches(8.9), Inches(1.9), Inches(3.9), Inches(0.06), ACCENT_GREEN)
tf3 = add_rich_textbox(slide, Inches(9.15), Inches(2.1), Inches(3.4), Inches(2.2))
add_para(tf3, "Q3：模型预测可靠性", font_size=17, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf3, "Ridge vs RF 时间切分对比 + 5-fold CV 验证", font_size=13, color=BLACK, space_after=6)
add_para(tf3, "残差诊断 + 适用场景分析 + 模型局限说明", font_size=13, color=GRAY, space_after=6)
add_para(tf3, "✅ 双模型对比分析 + 200+ 字分析", font_size=12, color=ACCENT_GREEN, space_after=4)

# 底部说明
add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.0), LIGHT_GRAY)
add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.06), ACCENT_BLUE)
tf4 = add_rich_textbox(slide, Inches(0.9), Inches(4.95), Inches(11.5), Inches(1.6))
add_para(tf4, "📋  设计思路", font_size=18, color=DARK_BLUE, bold=True, space_after=10)
add_para(tf4, "每道题对应一个独立 Tab，采用「图表 + 数据分析 + st.write 200+ 字论述」双层结构", font_size=14, color=BLACK, space_after=6, level=1)
add_para(tf4, "数据仅使用 3 个东北部镇区（Punggol / Sengkang / Hougang），不引入外部成熟区数据", font_size=14, color=BLACK, space_after=6, level=1)
add_para(tf4, "技术实现：Haversine MRT 距离 / CAGR 计算 / RF+Ridge 双模型 / 残差诊断 / OneHot 编码", font_size=14, color=BLACK, space_after=6, level=1)


# ============================================================
# 第 11 页：策略验证 — 4 策略体系
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 11)
slide_title(slide, "策略验证 — 4 策略 Train/Test 回测", "Strategy Validation")

# 4 策略说明
make_table(slide, Inches(0.3), Inches(1.85),
           [Inches(2.0), Inches(3.0), Inches(3.5), Inches(3.5)],
           ["策略", "逻辑", "选择规则", "适用场景"],
           [
               ["镇区偏离回归", "买低于镇区均价的房产\n等待价格回归", "训练期价格偏离度最大", "逆向投资"],
               ["MRT 便利", "买靠近 MRT 的房产\n享受交通溢价", "训练期 MRT 距离最近", "自住通勤"],
               ["长剩余租约", "买租约长的保值房产\n减少折旧损失", "训练期剩余租约最长", "长期持有"],
               ["低总价入门", "买低总价小户型\n降低入场门槛", "训练期总价最低", "首次购房"],
           ])

# 验证说明
card = add_rect(slide, Inches(0.3), Inches(4.2), Inches(12.5), Inches(2.5), LIGHT_GRAY)
add_rect(slide, Inches(0.3), Inches(4.2), Inches(12.5), Inches(0.06), ACCENT_GREEN)
tf = add_rich_textbox(slide, Inches(0.7), Inches(4.4), Inches(11.7), Inches(2.2))
add_para(tf, "🧪  验证设计", font_size=18, color=DARK_BLUE, bold=True, space_after=8)
add_para(tf, "硬 2023/2024 时间切分（防作弊）：2020–2023 策略形成期，2024+ 验证期", font_size=14, color=BLACK, space_after=6, level=1)
add_para(tf, "7 项验证指标：总涨幅 · CAGR · 价格波动 CV · 年度涨幅标准差 · 最大回撤 · 成交量 CV · 预算适配性", font_size=14, color=BLACK, space_after=6, level=1)
add_para(tf, "输出：策略组 vs 基准组对比表格 + 季度均价走势折线图 + 关键指标差异柱状图", font_size=14, color=BLACK, space_after=6, level=1)
add_para(tf, "策略反思：st.write 200+ 字（选择理由 / 优劣势 / 市场变化影响 / 改进方向）", font_size=14, color=BLACK, space_after=6, level=1)


# ============================================================
# 第 12 页：技术亮点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 12)
slide_title(slide, "技术亮点", "Technical Highlights")

highlights = [
    ("01", "完整数据管线", "API → Parquet 缓存 → 自动清洗 →\n3-sigma / 分位数异常值处理", PRIMARY_BLUE),
    ("02", "双模型对比", "线性 Ridge（可解释）+ 集成 RF（高精度）\n互补分析，分房型评估 6 类别", LIGHT_BLUE),
    ("03", "预测区间估计", "RF 个体树分位数 → 90% 置信区间\n四级渐进式放宽相似成交匹配", ACCENT_GREEN),
    ("04", "四策略量化验证", "Train/Test 时间切分回测\n7 项指标 + 策略反思 200+ 字", ACCENT_ORANGE),
    ("05", "Town vs Price 偏离度", "三镇区均价偏离度 + CAGR 双重评估\n替代传统成熟/非成熟二分法", ACCENT_RED),
    ("06", "一键部署", "内置 Parquet 缓存数据\nStreamlit Community Cloud 免费部署", DARK_BLUE),
]

for i, (num, title, desc, color) in enumerate(highlights):
    row = i // 3
    col = i % 3
    x = Inches(0.6) + col * Inches(4.15)
    y = Inches(2.0) + row * Inches(2.7)

    add_rect(slide, x, y, Inches(3.95), Inches(2.45), LIGHT_GRAY)
    add_rect(slide, x, y, Inches(3.95), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.8), Inches(0.5),
                num, font_size=32, color=color, bold=True, font_name="Consolas")
    add_textbox(slide, x + Inches(1.0), y + Inches(0.25), Inches(2.7), Inches(0.4),
                title, font_size=18, color=DARK_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.25), y + Inches(0.85), Inches(3.45), Inches(1.4),
                desc, font_size=12, color=GRAY)


# ============================================================
# 第 13 页：要求完成情况
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 13)
slide_title(slide, "要求完成情况", "Requirements Checklist")

make_table(slide, Inches(0.3), Inches(1.85),
           [Inches(3.3), Inches(1.0), Inches(4.0), Inches(4.0)],
           ["要求", "状态", "说明", "备注"],
           [
               ["3 个相邻镇区", "✅", "Punggol / Sengkang / Hougang", "新加坡东北部走廊"],
               ["2020 年至今数据", "✅", "2020-2026，清洗后 34,914 条", "远超最低要求（~16 倍）"],
               ["≥2,000 条记录", "✅", "34,914 条（16 倍+）", "清洗后高质量数据"],
               ["Streamlit 交互系统", "✅", "6 页面完整 Web 应用", "Cloud 一键部署"],
               ["筛选 + 数据表格", "✅", "9 项筛选 + CSV 导出", "可排序可配置列"],
               ["统计看板（含单价）", "✅", "8 指标卡片", "含平均/最高/最低单价"],
               ["地图 ×3 + 设施", "✅", "均价/热力/时空 3 模式", "MRT/学校叠加图层已实现"],
               ["≥4 个因素分析", "✅", "6 因素 Tab + 策略验证", "均含统计显著性检验"],
               ["≥2 个预测模型", "✅", "Ridge + Random Forest", "5-fold CV + 6 类别分房型评估"],
               ["购房策略 + 验证", "✅", "4 策略 Train/Test 回测", "7 项指标 + 策略反思"],
               ["分析思考题", "✅", "3 道题各 200+ 字分析", "保值性/策略/模型可靠性"],
           ])


# ============================================================
# 第 14 页：挑战与解决
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 14)
slide_title(slide, "项目挑战与解决方案", "Challenges & Solutions")

make_table(slide, Inches(0.3), Inches(1.9),
           [Inches(3.3), Inches(4.5), Inches(4.5)],
           ["挑战", "影响", "解决方案"],
           [
               ["CKAN API 分页上限", "~5 万条后截断", "按镇区分列请求，绕过全局限制"],
               ["时间切分下模型选择", "RF 无法外推时间趋势", "Ridge year 特征线性外推，时间切分评估"],
               ["相似成交匹配过严", "4 重过滤几乎无结果", "四级渐进放宽 + 兜底同镇区+同房型"],
               ["数据质量问题", "缺失/异常值干扰分析", "关键字段行删除 + 3-sigma 分组剔除"],
               ["热力图无真实坐标", "无法精确定位", "镇区聚类中心 + 随机抖动（~300m）"],
               ["Cloud 无法访问 API", "部署后数据加载失败", "内置 Parquet 缓存 + API 降级策略"],
               ["面积 × 房型共线性", "Ridge 系数不稳定", "L2 正则化 + RF 重要性聚合分析"],
               ["多页面代码组织", "维护困难", "pages/ 目录 + importlib 动态导入"],
           ])


# ============================================================
# 第 15 页：不足与改进
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 15)
slide_title(slide, "不足与改进方向", "Limitations & Future Work")

tf = add_rich_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5))
add_para(tf, "🔧  当前不足", font_size=22, color=DARK_BLUE, bold=True, space_after=14)
add_para(tf, "📍  地理位置为仿真坐标 → 可接入 OneMap API 获取精确经纬度", font_size=15, color=BLACK, space_after=8, level=1)
add_para(tf, "🤖  可增加 XGBoost / LightGBM / 时间序列模型（Prophet / LSTM）", font_size=15, color=BLACK, space_after=8, level=1)
add_para(tf, "🌍  宏观因素未纳入：利率、CPI、BTO 供应量、政策冲击", font_size=15, color=BLACK, space_after=8, level=1)
add_para(tf, "📐  政策量化待深入：RDD（断点回归）/ 事件研究法", font_size=15, color=BLACK, space_after=8, level=1)
add_para(tf, "🏫  配套设施数据仅 3 镇区 MRT 距离点，尚不足以独立建模分析", font_size=15, color=BLACK, space_after=8, level=1)

add_para(tf, "", font_size=8, color=BLACK, space_after=8)
add_para(tf, "🚀  未来改进方向", font_size=22, color=DARK_BLUE, bold=True, space_after=14)
add_para(tf, "接入真实地理坐标 → 精确空间分析（Moran's I 空间自相关）", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "增加 XGBoost / LightGBM 提升预测精度", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "纳入宏观经济指标，提升模型在拐点期的预测能力", font_size=15, color=BLACK, space_after=6, level=1)
add_para(tf, "政策影响量化：BSD/ABSD 调整前后的价格断点分析", font_size=15, color=BLACK, space_after=6, level=1)

# 免责声明
add_rect(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5), LIGHT_GRAY)
add_textbox(slide, Inches(1.0), Inches(6.65), Inches(11.1), Inches(0.4),
            "⚠️  免责声明：房价预测结果仅用于课程学习，不构成购房建议。",
            font_size=12, color=ACCENT_RED, alignment=CENTER)


# ============================================================
# 第 16 页：成员分工
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_bottom_bar(slide)
add_page_number(slide, 16)
slide_title(slide, "成员分工", "Team Contributions")

make_table(slide, Inches(0.8), Inches(2.2),
           [Inches(1.8), Inches(4.5), Inches(4.5), Inches(1.2)],
           ["成员", "负责工作", "具体内容", "占比"],
           [
               ["张涛",
                "数据管线 + 系统架构\n+ 数据概览 + 地图可视\n化 + 因素分析 + 分析思考题 + 部署",
                "data_loader.py / helpers.py / app.py\n01_overview.py / 02_map.py\n03_factors.py / 05_questions.py\n+ Cloud 部署 + 项目报告",
                "~50%"],
               ["申经纬",
                "价格预测 + 策略验证\n+ 答辩 PPT",
                "04_prediction.py（Ridge+RF 训练/评估\n/CV/特征重要性/交互式预测/相似成交）\n06_strategy.py（4 策略/回测/验证/反思）\n+ 答辩 PPT 制作",
                "~50%"],
           ])


# ============================================================
# 第 17 页：致谢 & Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), ACCENT_BLUE)

add_textbox(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(1.2),
            "感谢聆听，欢迎提问！", font_size=52, color=WHITE, bold=True, alignment=CENTER)

add_rect(slide, Inches(5.5), Inches(3.4), Inches(2.3), Inches(0.04), ACCENT_BLUE)

add_textbox(slide, Inches(0), Inches(3.8), Inches(13.333), Inches(0.6),
            "🌐  Live Demo", font_size=20, color=ACCENT_BLUE, alignment=CENTER)
add_textbox(slide, Inches(0), Inches(4.3), Inches(13.333), Inches(0.5),
            "https://hdb-resale-analysis-mxz6kgxpljqn4ubkrqrvr5.streamlit.app/",
            font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), alignment=CENTER, font_name="Consolas")

add_textbox(slide, Inches(0), Inches(5.0), Inches(13.333), Inches(0.5),
            "💻  本地运行：streamlit run app.py",
            font_size=16, color=RGBColor(0x94, 0xA3, 0xB8), alignment=CENTER, font_name="Consolas")

add_textbox(slide, Inches(0), Inches(5.8), Inches(13.333), Inches(0.6),
            "数据来源：data.gov.sg · HDB 官网 · IRAS  |  技术：Streamlit · Scikit-learn · Plotly · pydeck",
            font_size=12, color=GRAY, alignment=CENTER)


# ============================================================
# 保存
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "HDB答辩PPT.pptx")
prs.save(output_path)
print(f"PPT generated: {output_path}")
print(f"Total slides: {len(prs.slides)}")
